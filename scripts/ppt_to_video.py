"""PPT 转视频脚本 — 使用 python-pptx + Pillow 渲染幻灯片，imageio + ffmpeg 生成 MP4"""
import sys
import os
from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

import imageio
import numpy as np

# === 配置 ===
PPTX_PATH = r"e:\ap\AgentPrimordia\docs\AgentPrimordia-Promo.pptx"
OUTPUT_PATH = r"e:\ap\AgentPrimordia\docs\AgentPrimordia-Promo.mp4"
WIDTH = 1920   # 输出视频宽度
HEIGHT = 1080  # 输出视频高度
FPS = 30       # 帧率
SECONDS_PER_SLIDE = 5  # 每页停留秒数
FADE_FRAMES = 10       # 过渡淡入淡出帧数

# === EMU 转像素 ===
def emu_to_px(emu, scale):
    """将 EMU 单位转换为像素"""
    return int(emu * scale)

def get_font(size_pt, bold=False):
    """获取字体 — 尝试微软雅黑，回退到默认字体"""
    font_names = ["msyhbd.ttc" if bold else "msyh.ttc",
                  "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
                  "C:/Windows/Fonts/simhei.ttf"]
    for name in font_names:
        try:
            return ImageFont.truetype(name, size_pt)
        except (IOError, OSError):
            continue
    return ImageFont.load_default()

def rgb_from_element(elem):
    """从 XML 元素中提取 RGB 颜色"""
    # 尝试 srgbClr
    srgb = elem.find(qn('a:srgbClr'))
    if srgb is not None:
        val = srgb.get('val', '000000')
        return tuple(int(val[i:i+2], 16) for i in (0, 2, 4))
    # 尝试 schemeClr（简单映射）
    scheme = elem.find(qn('a:schemeClr'))
    if scheme is not None:
        # 默认返回深灰色
        return (30, 37, 60)
    return None

def get_fill_color(shape):
    """获取形状的填充颜色"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = sp.find(qn('a:spPr'))
    if spPr is None:
        return None
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        return rgb_from_element(solidFill)
    # 检查 noFill
    noFill = spPr.find(qn('a:noFill'))
    if noFill is not None:
        return None
    return None

def get_line_color(shape):
    """获取形状的线条颜色"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return None
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return None
    solidFill = ln.find(qn('a:solidFill'))
    if solidFill is not None:
        return rgb_from_element(solidFill)
    return None

def get_shape_alpha(shape):
    """获取形状填充的透明度"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return 255
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.find(qn('a:alpha'))
            if alpha_elem is not None:
                # alpha val 是百分比 * 1000，如 40000 = 40%
                val = int(alpha_elem.get('val', '100000'))
                return int(255 * val / 100000)
    return 255

def render_background(slide, img, scale_x, scale_y):
    """渲染幻灯片背景"""
    bg = slide.background
    if bg.fill.type is not None:
        try:
            color = bg.fill.fore_color.rgb
            r, g, b = int(str(color)[:2], 16), int(str(color)[2:4], 16), int(str(color)[4:6], 16)
            ImageDraw.Draw(img).rectangle([0, 0, WIDTH, HEIGHT], fill=(r, g, b))
        except Exception:
            pass

def safe_rgb(font_obj):
    """安全获取字体颜色 RGB，处理 NoneColor 等异常"""
    try:
        if font_obj.color and font_obj.color.type is not None and font_obj.color.rgb:
            c = str(font_obj.color.rgb)
            return tuple(int(c[i:i+2], 16) for i in (0, 2, 4))
    except (AttributeError, TypeError):
        pass
    return None

def render_text_frame(tf, left, top, width, height, draw, scale_x, scale_y):
    """渲染文本框内容"""
    x = emu_to_px(left, scale_x)
    y = emu_to_px(top, scale_y)
    w = emu_to_px(width, scale_x)
    h = emu_to_px(height, scale_y)

    cur_y = y
    for para in tf.paragraphs:
        if not para.text.strip():
            cur_y += 8
            continue
        # 获取段落字体属性
        font_size = 14
        font_color = (255, 255, 255)
        bold = False

        # 段落级别属性
        try:
            if para.font.size:
                font_size = int(para.font.size.pt * 1.33)
            fc = safe_rgb(para.font)
            if fc:
                font_color = fc
            if para.font.bold:
                bold = True
        except (AttributeError, TypeError):
            pass

        # run 级别属性（优先）
        for run in para.runs:
            try:
                if run.font.size:
                    font_size = int(run.font.size.pt * 1.33)
                fc = safe_rgb(run.font)
                if fc:
                    font_color = fc
                if run.font.bold is not None:
                    bold = run.font.bold
            except (AttributeError, TypeError):
                pass
            break  # 取第一个 run 的属性

        font = get_font(max(8, font_size), bold)
        text = para.text

        # 处理对齐
        alignment = para.alignment
        if alignment is not None and str(alignment) in ('CENTER', '1'):
            bbox = font.getbbox(text)
            tw = bbox[2] - bbox[0]
            text_x = x + (w - tw) // 2
        elif alignment is not None and str(alignment) in ('RIGHT', '2'):
            bbox = font.getbbox(text)
            tw = bbox[2] - bbox[0]
            text_x = x + w - tw
        else:
            text_x = x

        draw.text((text_x, cur_y), text, fill=font_color, font=font)
        cur_y += font_size + 4

        if cur_y > y + h:
            break

def render_picture(shape, draw, scale_x, scale_y):
    """渲染嵌入图片"""
    try:
        image_stream = shape.image.blob
        pic = Image.open(BytesIO(image_stream))
        left = shape.left
        top = shape.top
        w = shape.width if shape.width else pic.width * 914400 // 96  # 默认 96dpi 转 EMU
        h = shape.height if shape.height else pic.height * 914400 // 96

        px = emu_to_px(left, scale_x)
        py = emu_to_px(top, scale_y)
        pw = emu_to_px(w, scale_x)
        ph = emu_to_px(h, scale_y)

        if pw > 0 and ph > 0:
            pic_resized = pic.resize((pw, ph), Image.Resampling.LANCZOS)
            return pic_resized, px, py
    except Exception as e:
        pass
    return None, 0, 0

def render_table(shape, draw, scale_x, scale_y):
    """渲染表格"""
    table = shape.table
    left = shape.left
    top = shape.top
    width = shape.width
    row_count = len(table.rows)
    col_count = len(table.columns)

    total_height = emu_to_px(shape.height, scale_y) if shape.height else row_count * 30
    row_height = total_height // max(row_count, 1)

    cur_y = emu_to_px(top, scale_y)
    for ri, row in enumerate(table.rows):
        cur_x = emu_to_px(left, scale_x)
        for ci, cell in enumerate(row.cells):
            col_width = emu_to_px(table.columns[ci].width, scale_x)

            # 单元格背景
            cell_color = None
            try:
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr'))
                if tcPr is not None:
                    sf = tcPr.find(qn('a:solidFill'))
                    if sf is not None:
                        cell_color = rgb_from_element(sf)
            except Exception:
                pass

            if cell_color:
                draw.rectangle([cur_x, cur_y, cur_x + col_width, cur_y + row_height],
                             fill=cell_color)

            # 单元格文本
            text = cell.text.strip()
            if text:
                font_size = 13
                font_color = (226, 232, 240)  # LIGHT
                bold = False
                for p in cell.text_frame.paragraphs:
                    try:
                        if p.font.size:
                            font_size = int(p.font.size.pt * 1.33)
                        fc = safe_rgb(p.font)
                        if fc:
                            font_color = fc
                        if p.font.bold:
                            bold = True
                    except (AttributeError, TypeError):
                        pass
                    break
                font = get_font(max(8, font_size), bold)
                # 居中绘制
                bbox = font.getbbox(text)
                tw = bbox[2] - bbox[0]
                tx = cur_x + (col_width - tw) // 2
                ty = cur_y + (row_height - font_size) // 2
                draw.text((tx, ty), text, fill=font_color, font=font)

            cur_x += col_width
        cur_y += row_height

def render_slide(slide, slide_width, slide_height):
    """渲染单张幻灯片为 PIL Image — 两遍渲染: 先画形状背景, 再画文字"""
    scale_x = WIDTH / slide_width
    scale_y = HEIGHT / slide_height

    img = Image.new('RGB', (WIDTH, HEIGHT), (13, 17, 30))  # 默认深色背景
    draw = ImageDraw.Draw(img)

    # 渲染背景
    render_background(slide, img, scale_x, scale_y)
    draw = ImageDraw.Draw(img)  # 重新获取 draw 对象

    # 收集所有形状并按 z-order 渲染
    overlay_images = []  # 需要叠加的图片
    text_items = []      # 延迟渲染的文本 (shape, left, top, width, height)

    for shape in slide.shapes:
        left = shape.left if shape.left else 0
        top = shape.top if shape.top else 0
        width = shape.width if shape.width else 0
        height = shape.height if shape.height else 0

        # 图片
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_img, px, py = render_picture(shape, draw, scale_x, scale_y)
            if pic_img:
                overlay_images.append((pic_img, px, py))
            continue

        # 表格
        if shape.has_table:
            render_table(shape, draw, scale_x, scale_y)
            continue

        # 文本框 — 第一遍只画背景, 文字延迟到第二遍
        if shape.has_text_frame:
            fill_color = get_fill_color(shape)
            alpha = get_shape_alpha(shape)
            if fill_color:
                x = emu_to_px(left, scale_x)
                y = emu_to_px(top, scale_y)
                w = emu_to_px(width, scale_x)
                h = emu_to_px(height, scale_y)
                # 检查是否是圆角矩形
                is_rounded = False
                prstGeom = shape._element.find(qn('p:spPr'))
                if prstGeom is not None:
                    geom = prstGeom.find(qn('a:prstGeom'))
                    if geom is not None and geom.get('prst') == 'roundRect':
                        is_rounded = True
                # 绘制带透明度的矩形
                if alpha < 255:
                    overlay = Image.new('RGBA', (WIDTH, HEIGHT), (0, 0, 0, 0))
                    od = ImageDraw.Draw(overlay)
                    if is_rounded:
                        od.rounded_rectangle([x, y, x+w, y+h], radius=min(8, w//10, h//10),
                                            fill=(*fill_color, alpha))
                    else:
                        od.rectangle([x, y, x+w, y+h], fill=(*fill_color, alpha))
                    img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
                    draw = ImageDraw.Draw(img)
                else:
                    if is_rounded:
                        draw.rounded_rectangle([x, y, x+w, y+h],
                                             radius=min(8, w//10, h//10),
                                             fill=fill_color)
                    else:
                        draw.rectangle([x, y, x+w, y+h], fill=fill_color)

                # 绘制边框
                line_color = get_line_color(shape)
                if line_color:
                    if is_rounded:
                        draw.rounded_rectangle([x, y, x+w, y+h],
                                             radius=min(8, w//10, h//10),
                                             outline=line_color, width=1)
                    else:
                        draw.rectangle([x, y, x+w, y+h], outline=line_color, width=1)

            # 收集文本, 延迟到第二遍渲染
            text_items.append((shape.text_frame, left, top, width, height))
            continue

        # 纯形状（无文本框）— 绘制填充矩形
        fill_color = get_fill_color(shape)
        if fill_color:
            x = emu_to_px(left, scale_x)
            y = emu_to_px(top, scale_y)
            w = emu_to_px(width, scale_x)
            h = emu_to_px(height, scale_y)
            alpha = get_shape_alpha(shape)

            # 检查形状类型
            is_rounded = False
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                geom = spPr.find(qn('a:prstGeom'))
                if geom is not None and geom.get('prst') == 'roundRect':
                    is_rounded = True

            if alpha < 255:
                overlay = Image.new('RGBA', (WIDTH, HEIGHT), (0, 0, 0, 0))
                od = ImageDraw.Draw(overlay)
                if is_rounded:
                    od.rounded_rectangle([x, y, x+w, y+h], radius=min(8, w//10, h//10),
                                        fill=(*fill_color, alpha))
                else:
                    od.rectangle([x, y, x+w, y+h], fill=(*fill_color, alpha))
                img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
                draw = ImageDraw.Draw(img)
            else:
                if is_rounded:
                    draw.rounded_rectangle([x, y, x+w, y+h], radius=min(8, w//10, h//10),
                                          fill=fill_color)
                else:
                    draw.rectangle([x, y, x+w, y+h], fill=fill_color)

            # 绘制边框
            line_color = get_line_color(shape)
            if line_color:
                if is_rounded:
                    draw.rounded_rectangle([x, y, x+w, y+h], radius=min(8, w//10, h//10),
                                          outline=line_color, width=1)
                else:
                    draw.rectangle([x, y, x+w, y+h], outline=line_color, width=1)

    # 叠加图片
    for pic_img, px, py in overlay_images:
        if pic_img.mode != 'RGBA':
            pic_img = pic_img.convert('RGBA')
        img.paste(pic_img, (px, py), pic_img if pic_img.mode == 'RGBA' else None)

    # 第二遍: 在所有形状背景之上渲染文字
    draw = ImageDraw.Draw(img)
    for tf, left, top, width, height in text_items:
        render_text_frame(tf, left, top, width, height, draw, scale_x, scale_y)

    return img

def create_video(slides_images, output_path, fps=30, seconds_per_slide=5, fade_frames=10):
    """将幻灯片图片序列合成为 MP4 视频（带淡入淡出过渡）"""
    writer = imageio.get_writer(output_path, fps=fps, codec='libx264',
                                quality=8, pixelformat='yuv420p')

    total_slides = len(slides_images)
    hold_frames = seconds_per_slide * fps  # 每页停留帧数

    for i, img in enumerate(slides_images):
        frame = np.array(img.convert('RGB'))
        # 写入停留帧
        for _ in range(hold_frames):
            writer.append_data(frame)

        # 淡入淡出过渡到下一张
        if i < total_slides - 1:
            next_img = np.array(slides_images[i + 1].convert('RGB'))
            for f in range(fade_frames):
                alpha = (f + 1) / fade_frames
                blended = np.clip(
                    frame * (1 - alpha) + next_img * alpha,
                    0, 255
                ).astype(np.uint8)
                writer.append_data(blended)

    writer.close()
    print(f"视频已生成: {output_path}")

def main():
    print(f"正在读取 PPT: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)
    slide_count = len(prs.slides)
    print(f"共 {slide_count} 张幻灯片")

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 渲染所有幻灯片
    slides_images = []
    for idx, slide in enumerate(prs.slides):
        print(f"  渲染第 {idx + 1}/{slide_count} 张...")
        img = render_slide(slide, slide_width, slide_height)
        slides_images.append(img)

    # 生成视频
    print(f"\n正在生成视频...")
    create_video(slides_images, OUTPUT_PATH, FPS, SECONDS_PER_SLIDE, FADE_FRAMES)

    # 验证输出
    if os.path.exists(OUTPUT_PATH):
        size_mb = os.path.getsize(OUTPUT_PATH) / (1024 * 1024)
        print(f"输出文件大小: {size_mb:.2f} MB")
        if size_mb > 1:
            print("[OK] 视频生成成功！")
        else:
            print("[WARN] 文件偏小，可能存在问题")
    else:
        print("[FAIL] 视频生成失败")
        sys.exit(1)

if __name__ == "__main__":
    main()
